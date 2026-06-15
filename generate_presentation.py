import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Set dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Color Palette
    c_dark_bg = RGBColor(15, 23, 42)      # Deep Slate / Navy
    c_light_bg = RGBColor(248, 250, 252)  # Cool off-white
    c_white = RGBColor(255, 255, 255)
    c_primary_text = RGBColor(15, 23, 42) # Near black
    c_muted_text = RGBColor(71, 85, 105)   # Muted gray
    c_accent_teal = RGBColor(13, 148, 136) # Primary Teal
    c_accent_blue = RGBColor(37, 99, 235) # Secondary Blue
    c_accent_red = RGBColor(220, 38, 38)   # Diagnostic Red
    c_light_muted = RGBColor(148, 163, 184) # Light slide secondary text
    c_card_border = RGBColor(226, 232, 240) # Border line for cards
    
    # Font Settings
    f_sans = "Segoe UI"
    
    # Helper: Set slide solid background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Add title and subtitle for standard (light) slides
    def add_slide_header(slide, title_text, subtitle_text):
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = f_sans
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_primary_text
        
        # Subtitle text box
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = f_sans
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = c_accent_teal
        p_sub.font.bold = True
        
        # Horizontal Divider Line (Teal)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.65), Inches(11.733), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = c_accent_teal
        line.line.color.rgb = c_accent_teal

    # Helper: Add a text box paragraph
    def add_p(tf, text, font_size=13, font_color=c_primary_text, bold=False, space_after=6, italic=False):
        if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = f_sans
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.italic = italic
        p.space_after = Pt(space_after)
        return p

    # Helper: Add standard card backgrounds
    def draw_card(slide, left, top, width, height, bg_color=c_white, border_color=c_card_border):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, c_dark_bg)
    
    # Left accent block
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_accent_teal
    accent_bar.line.color.rgb = c_accent_teal
    
    # Title box (large)
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.6))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "🐳 AI Docker Natural Language Health Dashboard"
    p1.font.name = f_sans
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = c_white
    
    # Subtitle box
    sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(11.0), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "A Production-Ready AI Agent for Container Orchestration and Observability"
    p2.font.name = f_sans
    p2.font.size = Pt(18)
    p2.font.color.rgb = c_accent_teal
    p2.font.bold = True
    
    # Presenter Metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(11.0), Inches(1.0))
    tf_meta = meta_box.text_frame
    add_p(tf_meta, "Presented by: Antigravity AI Assistant", font_size=13, font_color=c_white, bold=True)
    add_p(tf_meta, "Technologies: Python | Streamlit | Docker SDK | Ollama | Plotly | SQLite", font_size=11, font_color=c_light_muted)
    add_p(tf_meta, "Version 1.0.0 (Production Ready)", font_size=10, font_color=c_accent_teal, italic=True)

    # ==========================================
    # SLIDE 2: Table of Contents (Light Theme)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, c_light_bg)
    add_slide_header(slide2, "Table of Contents", "Structure of the Presentation")
    
    # We will use 4 cards in a 2x2 grid to represent sections
    # Card layout details: Left / Top / Width / Height
    card_w = Inches(5.6)
    card_h = Inches(2.1)
    
    # Row 1 Left
    draw_card(slide2, Inches(0.8), Inches(2.1), card_w, card_h)
    tx = slide2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Part 1: Context & Core Value", font_size=15, font_color=c_accent_blue, bold=True, space_after=4)
    add_p(tf, "• Slide 3: The Problem Statement (Container Complexity)", font_size=12, space_after=2)
    add_p(tf, "• Slide 4: The Solution (Natural Language Ops)", font_size=12, space_after=2)
    add_p(tf, "• Slide 5: Project Objectives & System Design Pillars", font_size=12, space_after=2)
    
    # Row 1 Right
    draw_card(slide2, Inches(6.9), Inches(2.1), card_w, card_h)
    tx = slide2.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.2), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Part 2: Under the Hood Architecture", font_size=15, font_color=c_accent_blue, bold=True, space_after=4)
    add_p(tf, "• Slide 6-8: System Architecture & ReAct Loop Mechanics", font_size=12, space_after=2)
    add_p(tf, "• Slide 9: LLM Orchestration & llm.py Details", font_size=12, space_after=2)
    add_p(tf, "• Slide 10: Docker SDK Wrapper (docker_manager.py)", font_size=12, space_after=2)

    # Row 2 Left
    draw_card(slide2, Inches(0.8), Inches(4.6), card_w, card_h)
    tx = slide2.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(5.2), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Part 3: UI Pages & Features", font_size=15, font_color=c_accent_blue, bold=True, space_after=4)
    add_p(tf, "• Slide 11-13: UI Overview, Dashboard & Containers grid", font_size=12, space_after=2)
    add_p(tf, "• Slide 14: The AI Agent Conversational Page", font_size=12, space_after=2)
    add_p(tf, "• Slide 15: Telemetry, Logs, & External GitHub APIs", font_size=12, space_after=2)

    # Row 2 Right
    draw_card(slide2, Inches(6.9), Inches(4.6), card_w, card_h)
    tx = slide2.shapes.add_textbox(Inches(7.1), Inches(4.7), Inches(5.2), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Part 4: Operations, Deployment & Future", font_size=15, font_color=c_accent_blue, bold=True, space_after=4)
    add_p(tf, "• Slide 16-17: Database Schemas & Production Deployment", font_size=12, space_after=2)
    add_p(tf, "• Slide 18-19: Troubleshooting, Testing & QA Checklists", font_size=12, space_after=2)
    add_p(tf, "• Slide 20: Future Scope, Roadmap & Conclusion", font_size=12, space_after=2)

    # ==========================================
    # SLIDE 3: The Problem Statement (Light Theme)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, c_light_bg)
    add_slide_header(slide3, "The Problem: Container Complexity", "Why Direct Container Administration is Fragile")
    
    card_w = Inches(3.64)
    card_h = Inches(4.6)
    
    # Card 1: CLI Barriers
    draw_card(slide3, Inches(0.8), Inches(2.1), card_w, card_h)
    tx = slide3.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "1. Command Line Barriers", font_size=16, font_color=c_accent_red, bold=True, space_after=12)
    add_p(tf, "• Syntax Complexity", font_size=13, bold=True, space_after=4)
    add_p(tf, "Docker CLI flags (-d, -p, -v, --network, --restart) require exact syntax rules.", font_size=11, font_color=c_muted_text, space_after=8)
    add_p(tf, "• Risk of Mistakes", font_size=13, bold=True, space_after=4)
    add_p(tf, "A single typo during a critical production incident can shut down core applications.", font_size=11, font_color=c_muted_text, space_after=8)
    add_p(tf, "• Training Overhead", font_size=13, bold=True, space_after=4)
    add_p(tf, "Requires junior operators to memorize hundreds of command schemas.", font_size=11, font_color=c_muted_text)

    # Card 2: Fragmented Observability
    draw_card(slide3, Inches(4.84), Inches(2.1), card_w, card_h)
    tx = slide3.shapes.add_textbox(Inches(5.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "2. Fragmented Telemetry", font_size=16, font_color=c_accent_red, bold=True, space_after=12)
    add_p(tf, "• Isolated Streams", font_size=13, bold=True, space_after=4)
    add_p(tf, "Container lists, stats, text logs, and audit logs reside in completely disconnected interfaces.", font_size=11, font_color=c_muted_text, space_after=8)
    add_p(tf, "• No Historical Logging", font_size=13, bold=True, space_after=4)
    add_p(tf, "Docker doesn't save CPU/Memory trends locally without setting up complex external stacks (Prometheus).", font_size=11, font_color=c_muted_text, space_after=8)
    add_p(tf, "• Lack of Audit Trails", font_size=13, bold=True, space_after=4)
    add_p(tf, "Identifying who stopped a container or why a container restarted is difficult without unified audits.", font_size=11, font_color=c_muted_text)

    # Card 3: Operation Bottlenecks
    draw_card(slide3, Inches(8.88), Inches(2.1), card_w, card_h)
    tx = slide3.shapes.add_textbox(Inches(9.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "3. Team Bottlenecks", font_size=16, font_color=c_accent_red, bold=True, space_after=12)
    add_p(tf, "• Operations Locks", font_size=13, bold=True, space_after=4)
    add_p(tf, "Non-technical stakeholders are blocked from managing test environments.", font_size=11, font_color=c_muted_text, space_after=8)
    add_p(tf, "• CLI Gatekeeping", font_size=13, bold=True, space_after=4)
    add_p(tf, "Production credentials are only shared with senior staff, creating deployment bottlenecks.", font_size=11, font_color=c_muted_text, space_after=8)
    add_p(tf, "• slow MTTR", font_size=13, bold=True, space_after=4)
    add_p(tf, "Locating a crashed container requires logging into VMs and parsing text streams.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 4: The Solution (Light Theme)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, c_light_bg)
    add_slide_header(slide4, "The Solution: Natural Language Operations", "Simplifying Container Management with AI and Telemetry")
    
    # Left Description Column
    tx_desc = slide4.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.6))
    tf_desc = tx_desc.text_frame
    tf_desc.word_wrap = True
    add_p(tf_desc, "Natural Language Interface", font_size=18, font_color=c_accent_teal, bold=True, space_after=10)
    add_p(tf_desc, "By mapping human statements directly to Docker API calls, the system simplifies systems administration. Users don't need to know complex flags; they describe their target state in plain English.", font_size=13, space_after=15)
    add_p(tf_desc, "Key Operational Enhancements:", font_size=14, bold=True, space_after=8)
    add_p(tf_desc, "• Lower Cognitive Load: Describe intentions without technical jargon.", font_size=12, space_after=4)
    add_p(tf_desc, "• ReAct Safety Loop: The AI reasons and reviews tool parameters before hitting the socket API.", font_size=12, space_after=4)
    add_p(tf_desc, "• Local Privacy: Running Ollama locally ensures zero external network exposure for server command prompts.", font_size=12, space_after=4)
    add_p(tf_desc, "• Audit Trail Security: Every prompt and action is captured.", font_size=12)

    # Right Card: Before vs. After comparison
    draw_card(slide4, Inches(6.8), Inches(2.1), Inches(5.7), Inches(4.6))
    tx_comp = slide4.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_comp = tx_comp.text_frame
    tf_comp.word_wrap = True
    add_p(tf_comp, "Operational Comparison", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf_comp, "TRADITIONAL CLI APPROACH (Manual & Fragile)", font_size=12, bold=True, space_after=4)
    add_p(tf_comp, "User inputs complex command in terminal:", font_size=11, space_after=2)
    add_p(tf_comp, "$ docker restart $(docker ps -aq --filter name=mysql)", font_size=11, font_color=c_accent_red, bold=True, space_after=8, italic=True)
    add_p(tf_comp, "• Risks: typos, wrong container matching, missing flags.", font_size=10, font_color=c_muted_text, space_after=12)
    
    add_p(tf_comp, "NATURAL LANGUAGE APPROACH (Declarative & Safe)", font_size=12, bold=True, space_after=4)
    add_p(tf_comp, "User inputs plain English instruction in chat box:", font_size=11, space_after=2)
    add_p(tf_comp, "\"restart mysql-test\"", font_size=12, font_color=c_accent_teal, bold=True, space_after=8, italic=True)
    add_p(tf_comp, "• AI handles mapping: Resolves image name, matches docker socket ID, triggers API restart action, and reports a readable success message.", font_size=10, font_color=c_muted_text)

    # ==========================================
    # SLIDE 5: Core Project Objectives (Light Theme)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, c_light_bg)
    add_slide_header(slide5, "Core Project Objectives", "Architectural Pillars of the Application Design")
    
    card_w = Inches(2.7)
    card_h = Inches(4.6)
    
    # 4 columns grid
    # Card 1: Accurate NLP Translation
    draw_card(slide5, Inches(0.8), Inches(2.1), card_w, card_h)
    tx = slide5.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(2.4), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "1. NLP Mapping", font_size=15, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Translate statements into strict JSON structures. Ensures high predictability and prevents code injections. Handles typos and slang gracefully.", font_size=11, font_color=c_muted_text)

    # Card 2: Step-by-Step Reason
    draw_card(slide5, Inches(3.8), Inches(2.1), card_w, card_h)
    tx = slide5.shapes.add_textbox(Inches(3.95), Inches(2.3), Inches(2.4), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "2. ReAct Trace", font_size=15, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Split LLM logic into reasoning steps. Make choices visible in the UI. Operator reviews action parameters before commands execute.", font_size=11, font_color=c_muted_text)

    # Card 3: Visual Analytics
    draw_card(slide5, Inches(6.8), Inches(2.1), card_w, card_h)
    tx = slide5.shapes.add_textbox(Inches(6.95), Inches(2.3), Inches(2.4), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "3. Visual Metrics", font_size=15, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Display CPU, Memory, and Status changes over time. Uses interactive Plotly charts. Prevents command line tracking bottlenecks.", font_size=11, font_color=c_muted_text)

    # Card 4: Local Isolation
    draw_card(slide5, Inches(9.8), Inches(2.1), card_w, card_h)
    tx = slide5.shapes.add_textbox(Inches(9.95), Inches(2.3), Inches(2.4), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "4. Production Ready", font_size=15, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Uses a real SQLite backend, live Docker socket integrations, and local LLM runtime. Zero mocked endpoints or placeholder layouts.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 6: System Architecture Divider (Dark Theme)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, c_dark_bg)
    
    accent_bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_accent_teal
    accent_bar.line.color.rgb = c_accent_teal
    
    title_box = slide6.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.6))
    tf6 = title_box.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "System Architecture & Data Flow"
    p6.font.name = f_sans
    p6.font.size = Pt(36)
    p6.font.bold = True
    p6.font.color.rgb = c_white
    
    sub_box = slide6.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(11.0), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Part 2: Understanding the Core Code Modules and Orchestration Pipeline"
    p2.font.name = f_sans
    p2.font.size = Pt(18)
    p2.font.color.rgb = c_accent_teal
    p2.font.bold = True

    # ==========================================
    # SLIDE 7: High-Level System Architecture (Light Theme)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, c_light_bg)
    add_slide_header(slide7, "High-Level System Architecture", "Data and Communication Pipelines")
    
    # Left Column: Structure
    draw_card(slide7, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide7.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Architecture Layers & Files", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf, "1. USER INTERFACE LAYER (app.py)", font_size=12, bold=True, space_after=2)
    add_p(tf, "Multi-page Streamlit dashboard providing cards, grids, Plotly charts, and chat interfaces.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "2. AI ORCHESTRATION LAYER (agent.py, llm.py)", font_size=12, bold=True, space_after=2)
    add_p(tf, "Takes command string, invokes local Ollama translator, processes output schema, and triggers tools.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "3. SYSTEM BOUNDARY LAYER (docker_manager.py, github_api.py)", font_size=12, bold=True, space_after=2)
    add_p(tf, "Speaks to local Docker UNIX Socket API and external GitHub Status endpoint.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "4. PERSISTENCE LAYER (database.py)", font_size=12, bold=True, space_after=2)
    add_p(tf, "Writes prompts, execution latency, metric snapshot records to SQLite.", font_size=11, font_color=c_muted_text)

    # Right Column: Visual Pipeline
    draw_card(slide7, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_vis = slide7.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_vis = tx_vis.text_frame
    tf_vis.word_wrap = True
    add_p(tf_vis, "Execution Data Pipeline", font_size=16, font_color=c_accent_blue, bold=True, space_after=15)
    
    # 4 horizontal visual steps
    add_p(tf_vis, "Step 1: Input Ingestion", font_size=13, bold=True, space_after=2)
    add_p(tf_vis, "User inputs command \"restart nginx\" in the Streamlit agent chat interface.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf_vis, "Step 2: LLM Interpretation & Mapping", font_size=13, bold=True, space_after=2)
    add_p(tf_vis, "llm.py calls Ollama to map input into a structured action: restart, target: nginx.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf_vis, "Step 3: Docker SDK Trigger", font_size=13, bold=True, space_after=2)
    add_p(tf_vis, "docker_manager.py binds to /var/run/docker.sock and restarts the container.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf_vis, "Step 4: SQL Logger", font_size=13, bold=True, space_after=2)
    add_p(tf_vis, "database.py records the command, result string, and execution duration in SQLite.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 8: The ReAct Agent Loop Concept (Light Theme)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, c_light_bg)
    add_slide_header(slide8, "The ReAct Agent Loop Paradigm", "Ensuring Traceable and Controllable AI Actions")
    
    # Left Column: Theory
    draw_card(slide8, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide8.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Why ReAct (Reasoning + Acting)?", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Standard LLM setups issue commands blindly in one step, masking internal steps. Our ReAct loop breaks down inputs into visible, audit-ready steps:", font_size=12, space_after=12)
    
    add_p(tf, "• Transparency", font_size=13, bold=True, space_after=2)
    add_p(tf, "Every decision, tool choice, and argument is printed dynamically in the UI.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Error Isolation", font_size=13, bold=True, space_after=2)
    add_p(tf, "If the LLM generates bad JSON or target matching fails, the step highlights the exact issue.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Permanent Tracing", font_size=13, bold=True, space_after=2)
    add_p(tf, "Every intermediate step is saved to the SQLite prompt log table for diagnostic audits.", font_size=11, font_color=c_muted_text)

    # Right Column: Visual execution flow representation
    draw_card(slide8, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_vis = slide8.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_vis = tx_vis.text_frame
    tf_vis.word_wrap = True
    add_p(tf_vis, "Orchestration Agent Loops", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    # Horizontal flow blocks
    add_p(tf_vis, "[User Request Ingested]", font_size=12, bold=True, space_after=2)
    add_p(tf_vis, "Initial user query is received and validated by the Streamlit state controller.", font_size=11, font_color=c_muted_text, space_after=10)
    
    add_p(tf_vis, "[Thought -> Action -> Observation]", font_size=12, bold=True, space_after=2)
    add_p(tf_vis, "• Thought: Decide which container API action maps to the user text.\n• Action: Translate query parameters and execute docker_manager tools.\n• Observation: Catch API feedback, logs, or diagnostic errors.", font_size=11, font_color=c_muted_text, space_after=10)
    
    add_p(tf_vis, "[Response Summary & Log Hydration]", font_size=12, bold=True, space_after=2)
    add_p(tf_vis, "Translate raw JSON SDK statistics back to natural English summaries while logging the execution steps.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 9: ReAct Loop Step-by-Step Details (Light Theme)
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, c_light_bg)
    add_slide_header(slide9, "ReAct Loop: 7 Steps of Execution", "Step-by-step trace of how the agent resolves commands")
    
    # We will draw 7 small horizontal/vertical blocks
    # 2 rows layout (4 on top, 3 on bottom)
    card_w = Inches(2.7)
    card_h = Inches(2.1)
    
    # Row 1
    # Step 1
    draw_card(slide9, Inches(0.8), Inches(2.1), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(0.95), Inches(2.2), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 1: User Request", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Capture user command in UI; filter out blank prompts and execute basic security guards.", font_size=11, font_color=c_muted_text)

    # Step 2
    draw_card(slide9, Inches(3.8), Inches(2.1), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(3.95), Inches(2.2), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 2: AI Analysis", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Submit user prompt and system constraints context to local Ollama (llama3) LLM.", font_size=11, font_color=c_muted_text)

    # Step 3
    draw_card(slide9, Inches(6.8), Inches(2.1), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(6.95), Inches(2.2), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 3: Tool Selection", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Extract action parameters from the locked JSON structure (e.g. action: restart, name: nginx).", font_size=11, font_color=c_muted_text)

    # Step 4
    draw_card(slide9, Inches(9.8), Inches(2.1), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(9.95), Inches(2.2), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 4: Tool Execution", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Map the parsed parameters directly to python docker SDK client methods inside docker_manager.", font_size=11, font_color=c_muted_text)

    # Row 2 (3 columns centered: left=2.3, space=0.3)
    # Step 5
    draw_card(slide9, Inches(2.3), Inches(4.6), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(2.45), Inches(4.7), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 5: Result Collection", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Collect execution stats, log records, and runtime duration metric details.", font_size=11, font_color=c_muted_text)

    # Step 6
    draw_card(slide9, Inches(5.3), Inches(4.6), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(5.45), Inches(4.7), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 6: Summary Design", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Feed action data back into LLM to output a human-friendly execution status recap.", font_size=11, font_color=c_muted_text)

    # Step 7
    draw_card(slide9, Inches(8.3), Inches(4.6), card_w, card_h)
    tx = slide9.shapes.add_textbox(Inches(8.45), Inches(4.7), Inches(2.4), Inches(1.9))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Step 7: Render & Log", font_size=13, font_color=c_accent_teal, bold=True, space_after=4)
    add_p(tf, "Render summaries and steps in UI; trigger parallel logging writes to SQLite.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 10: Core Component: llm.py (Light Theme)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, c_light_bg)
    add_slide_header(slide10, "Core Component: llm.py", "Local LLM Integration & Schema Enforcement")
    
    # Left Column: Configuration
    draw_card(slide10, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide10.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "LLM Orchestration Details", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf, "• Local-First API Connections", font_size=13, bold=True, space_after=2)
    add_p(tf, "Interfaces with Ollama container runtime (defaulting to llama3 model). Supports remote endpoints (like Groq) via API key settings overrides.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• High Predictability settings", font_size=13, bold=True, space_after=2)
    add_p(tf, "Uses temperature 0.1 to avoid creative variations and enforce deterministic responses.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Strict JSON Schema Locking", font_size=13, bold=True, space_after=2)
    add_p(tf, "System prompt overrides compel the LLM to output a clean JSON containing exact fields: 'action', 'container', 'status', etc. Prevents conversation text headers.", font_size=11, font_color=c_muted_text)

    # Right Column: The System Prompt Details
    draw_card(slide10, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_prompt = slide10.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_prompt = tx_prompt.text_frame
    tf_prompt.word_wrap = True
    add_p(tf_prompt, "System Prompt Schema Constraints", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf_prompt, "Few-shot templates in prompts/system_prompt.txt bind LLM outputs to structured maps:", font_size=12, space_after=8)
    
    # Code block format representation
    draw_card(slide10, Inches(7.1), Inches(3.2), Inches(5.2), Inches(3.2), bg_color=c_dark_bg, border_color=c_dark_bg)
    tx_code = slide10.shapes.add_textbox(Inches(7.2), Inches(3.3), Inches(5.0), Inches(3.0))
    tf_code = tx_code.text_frame
    tf_code.word_wrap = True
    add_p(tf_code, "{", font_size=11, font_color=c_accent_teal)
    add_p(tf_code, "  \"action\": \"restart_container\" | \"stop_container\" |", font_size=11, font_color=c_white)
    add_p(tf_code, "            \"start_container\" | \"list_containers\" |", font_size=11, font_color=c_white)
    add_p(tf_code, "            \"get_container_logs\" | \"get_container_health\",", font_size=11, font_color=c_white)
    add_p(tf_code, "  \"container\": \"nginx-test\" | \"redis-test\" | ... ,", font_size=11, font_color=c_white)
    add_p(tf_code, "  \"status\": \"running\" | \"all\" | \"exited\" | ...", font_size=11, font_color=c_white)
    add_p(tf_code, "}", font_size=11, font_color=c_accent_teal)

    # ==========================================
    # SLIDE 11: UI Pages Divider (Dark Theme)
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, c_dark_bg)
    
    accent_bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_accent_teal
    accent_bar.line.color.rgb = c_accent_teal
    
    title_box = slide11.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.6))
    tf11 = title_box.text_frame
    tf11.word_wrap = True
    p11 = tf11.paragraphs[0]
    p11.text = "User Interface & Streamlit Pages"
    p11.font.name = f_sans
    p11.font.size = Pt(36)
    p11.font.bold = True
    p11.font.color.rgb = c_white
    
    sub_box = slide11.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(11.0), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Part 3: Exploring the Multi-Page Streamlit Operations Dashboard"
    p2.font.name = f_sans
    p2.font.size = Pt(18)
    p2.font.color.rgb = c_accent_teal
    p2.font.bold = True

    # ==========================================
    # SLIDE 12: Core Component: docker_manager.py (Light Theme)
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, c_light_bg)
    add_slide_header(slide12, "Core Component: docker_manager.py", "Safe Programmatic Docker SDK Integration")
    
    # Left Column: Binding details
    draw_card(slide12, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide12.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Docker API Socket Binding", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf, "• Cross-Platform Socket Matching", font_size=13, bold=True, space_after=2)
    add_p(tf, "Initializes connection dynamically. Binds to UNIX socket /var/run/docker.sock on Linux/macOS. Falls back to npipe:////./pipe/docker_engine on Windows environments.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Safe Error Interception", font_size=13, bold=True, space_after=2)
    add_p(tf, "Catches connection and socket permission exceptions gracefully. Avoids dashboard crashes by returning a structured error JSON when Docker is offline.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Isolation Architecture", font_size=13, bold=True, space_after=2)
    add_p(tf, "Only allows read/write commands matching specific targets. Disallows execution of dangerous arbitrary script lines.", font_size=11, font_color=c_muted_text)

    # Right Column: Exposed API Actions
    draw_card(slide12, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_api = slide12.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_api = tx_api.text_frame
    tf_api.word_wrap = True
    add_p(tf_api, "Exposed API Actions", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    # List of SDK operations
    add_p(tf_api, "• list_containers(status)", font_size=12, bold=True, space_after=2)
    add_p(tf_api, "Queries container details using SDK filters: running, exited, restarting, or all.", font_size=10, font_color=c_muted_text, space_after=6)
    
    add_p(tf_api, "• start / stop / restart_container(name)", font_size=12, bold=True, space_after=2)
    add_p(tf_api, "Triggers container transitions. Uses try/except boundaries to intercept API errors.", font_size=10, font_color=c_muted_text, space_after=6)
    
    add_p(tf_api, "• get_container_logs(name, lines=100)", font_size=12, bold=True, space_after=2)
    add_p(tf_api, "Pulls raw container logs from standard output streams.", font_size=10, font_color=c_muted_text, space_after=6)
    
    add_p(tf_api, "• get_container_resource_usage(name)", font_size=12, bold=True, space_after=2)
    add_p(tf_api, "Calculates CPU and Memory utilization using stats deltas (system vs. container usage).", font_size=10, font_color=c_muted_text)

    # ==========================================
    # SLIDE 13: Pages 1 & 2: Dashboard & Containers (Light Theme)
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, c_light_bg)
    add_slide_header(slide13, "Pages 1 & 2: Dashboard & Containers", "Real-time Telemetry and Operational Controls")
    
    # Left Card: Dashboard Page
    draw_card(slide13, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide13.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Page 1: Telemetry Dashboard", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Provides rapid, high-level operational status statistics:", font_size=12, space_after=8)
    
    add_p(tf, "• Real-time KPI metrics cards", font_size=12, bold=True, space_after=2)
    add_p(tf, "Displays total containers count, running containers, stopped containers, restarting containers, and unhealthy container counts.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Diagnostic Quick Actions", font_size=12, bold=True, space_after=2)
    add_p(tf, "Shortcuts allow users to reload page details or trigger immediate system checks.", font_size=11, font_color=c_muted_text)

    # Right Card: Containers Page
    draw_card(slide13, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_grid = slide13.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_grid = tx_grid.text_frame
    tf_grid.word_wrap = True
    add_p(tf_grid, "Page 2: Container Management Grid", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf_grid, "Enables standard, click-to-run docker adjustments:", font_size=12, space_after=8)
    
    add_p(tf_grid, "• Clean Pandas DataFrame Table", font_size=12, bold=True, space_after=2)
    add_p(tf_grid, "Displays Name, Status, Image Tag, CPU%, Memory MB, and Container ID in a sortable, paginated grid.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf_grid, "• Inline Command Buttons", font_size=12, bold=True, space_after=2)
    add_p(tf_grid, "Allows administrators to click Start, Stop, or Restart buttons directly in the grid rows, bypassing CLI terminal login operations.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 14: Page 3: The AI Agent Interface (Light Theme)
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, c_light_bg)
    add_slide_header(slide14, "Page 3: AI Agent Chat Interface", "Conversational operations console with step-by-step audit tracing")
    
    # Left Column: Features
    draw_card(slide14, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide14.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Natural Language Console Features", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf, "• Interactive Chat Ingestion", font_size=13, bold=True, space_after=2)
    add_p(tf, "Simple prompt entry field. Recommends standard operations (e.g. \"show running containers\", \"restart nginx-test\").", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Real-Time Progress Indicators", font_size=13, bold=True, space_after=2)
    add_p(tf, "Displays step-by-step checklist ticks (✅/⏳/❌) as the agent progresses through the 7-step loop.", font_size=11, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "• Summary Result block", font_size=13, bold=True, space_after=2)
    add_p(tf, "Renders simple natural language execution recap along with the raw action JSON return.", font_size=11, font_color=c_muted_text)

    # Right Column: Visual Layout of the Agent screen
    draw_card(slide14, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_vis = slide14.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_vis = tx_vis.text_frame
    tf_vis.word_wrap = True
    add_p(tf_vis, "Visualizing Agent Execution Grids", font_size=16, font_color=c_accent_blue, bold=True, space_after=10)
    
    # Render simulated screen output
    draw_card(slide14, Inches(7.1), Inches(2.9), Inches(5.2), Inches(3.5), bg_color=c_dark_bg, border_color=c_dark_bg)
    tx_scr = slide14.shapes.add_textbox(Inches(7.2), Inches(3.0), Inches(5.0), Inches(3.3))
    tf_scr = tx_scr.text_frame
    tf_scr.word_wrap = True
    add_p(tf_scr, "User prompt: restart nginx-test", font_size=11, font_color=c_accent_teal, bold=True, space_after=6)
    add_p(tf_scr, "Checking system logs ...", font_size=10, font_color=c_white)
    add_p(tf_scr, "[ReAct Loop Active]", font_size=10, font_color=c_light_muted)
    add_p(tf_scr, " ✅ Step 1: User Request Ingested", font_size=10, font_color=c_white)
    add_p(tf_scr, " ✅ Step 2: AI Translation Finished", font_size=10, font_color=c_white)
    add_p(tf_scr, " ✅ Step 3: Tool Selected (restart_container)", font_size=10, font_color=c_white)
    add_p(tf_scr, " ✅ Step 4: Docker API Action Executed", font_size=10, font_color=c_white)
    add_p(tf_scr, " ... System Status Summary: nginx-test restarted successfully in 850ms.", font_size=10, font_color=c_accent_teal)

    # ==========================================
    # SLIDE 15: Pages 4, 5, 6 & 7: Analytics, Logs & Settings (Light Theme)
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, c_light_bg)
    add_slide_header(slide15, "Pages 4-7: Analytics, Logs, settings & GitHub Feed", "Diagnostics, rate limiting, and system configuration")
    
    card_w = Inches(3.64)
    card_h = Inches(4.6)
    
    # 3 columns layout
    # Page 4: Analytics
    draw_card(slide15, Inches(0.8), Inches(2.1), card_w, card_h)
    tx = slide15.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Page 4: Analytics", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Pulls database stats and plots metrics using Plotly:", font_size=11, space_after=8)
    add_p(tf, "• CPU & Memory lines", font_size=12, bold=True, space_after=2)
    add_p(tf, "Draws trend charts of resource usage snapshots.", font_size=11, font_color=c_muted_text, space_after=6)
    add_p(tf, "• Status Distributions", font_size=12, bold=True, space_after=2)
    add_p(tf, "Pie and bar charts dividing running, stopped, and crashed containers.", font_size=11, font_color=c_muted_text)

    # Page 5 & 6: Logs & GitHub
    draw_card(slide15, Inches(4.84), Inches(2.1), card_w, card_h)
    tx = slide15.shapes.add_textbox(Inches(5.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Page 5 & 6: Logs & GitHub", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "• Full Diagnostic Logs", font_size=12, bold=True, space_after=2)
    add_p(tf, "Inspect prompt history, database write lists, and container audit logs.", font_size=11, font_color=c_muted_text, space_after=6)
    add_p(tf, "• GitHub API Feed", font_size=12, bold=True, space_after=2)
    add_p(tf, "Tracks GitHub Status service health checks and pulls public events timeline under API rate limiting caps.", font_size=11, font_color=c_muted_text)

    # Page 7: Settings
    draw_card(slide15, Inches(8.88), Inches(2.1), card_w, card_h)
    tx = slide15.shapes.add_textbox(Inches(9.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Page 7: System Settings", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Configuration panel for connection settings:", font_size=11, space_after=8)
    add_p(tf, "• Docker Host override", font_size=12, bold=True, space_after=2)
    add_p(tf, "Configures socket paths or remote daemon ports.", font_size=11, font_color=c_muted_text, space_after=6)
    add_p(tf, "• LLM Server settings", font_size=12, bold=True, space_after=2)
    add_p(tf, "Adjusts Ollama endpoint addresses, model names, or API keys.", font_size=11, font_color=c_muted_text)

    # ==========================================
    # SLIDE 16: Ops, Setup & Security Divider (Dark Theme)
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, c_dark_bg)
    
    accent_bar = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_accent_teal
    accent_bar.line.color.rgb = c_accent_teal
    
    title_box = slide16.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(1.6))
    tf16 = title_box.text_frame
    tf16.word_wrap = True
    p16 = tf16.paragraphs[0]
    p16.text = "Deployment, Persistence & Quality QA"
    p16.font.name = f_sans
    p16.font.size = Pt(36)
    p16.font.bold = True
    p16.font.color.rgb = c_white
    
    sub_box = slide16.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(11.0), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Part 4: Database Schemas, Deployment Frameworks, and QA Protocols"
    p2.font.name = f_sans
    p2.font.size = Pt(18)
    p2.font.color.rgb = c_accent_teal
    p2.font.bold = True

    # ==========================================
    # SLIDE 17: Database Schemas & Persistence (Light Theme)
    # ==========================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, c_light_bg)
    add_slide_header(slide17, "Data Model: Database & File Schemas", "SQLite persistence tables mapping audit trails")
    
    card_w = Inches(3.64)
    card_h = Inches(4.6)
    
    # 3 columns for 3 tables
    # Table 1: prompt_logs
    draw_card(slide17, Inches(0.8), Inches(2.1), card_w, card_h)
    tx = slide17.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "1. Table: prompt_logs", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Logs user command strings and their corresponding LLM actions.", font_size=11, space_after=8)
    add_p(tf, "Key fields schema:", font_size=12, bold=True, space_after=4)
    add_p(tf, "• id (INTEGER PRIMARY KEY)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• timestamp (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• user_prompt (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• generated_action (TEXT)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• execution_result (TEXT)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• execution_time_ms (REAL)", font_size=10, font_color=c_muted_text)

    # Table 2: container_history
    draw_card(slide17, Inches(4.84), Inches(2.1), card_w, card_h)
    tx = slide17.shapes.add_textbox(Inches(5.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "2. Table: container_history", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Tracks telemetry resource usage metrics over time.", font_size=11, space_after=8)
    add_p(tf, "Key fields schema:", font_size=12, bold=True, space_after=4)
    add_p(tf, "• id (INTEGER PRIMARY KEY)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• timestamp (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• container_id (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• container_name (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• status (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• cpu_usage / memory_usage (REAL)", font_size=10, font_color=c_muted_text)

    # Table 3: audit_logs
    draw_card(slide17, Inches(8.88), Inches(2.1), card_w, card_h)
    tx = slide17.shapes.add_textbox(Inches(9.0), Inches(2.3), Inches(3.24), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "3. Table: audit_logs", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Maintains system auditing trace records.", font_size=11, space_after=8)
    add_p(tf, "Key fields schema:", font_size=12, bold=True, space_after=4)
    add_p(tf, "• id (INTEGER PRIMARY KEY)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• timestamp (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• action (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• status (TEXT NOT NULL)", font_size=10, font_color=c_muted_text, space_after=2)
    add_p(tf, "• details (TEXT)", font_size=10, font_color=c_muted_text)

    # ==========================================
    # SLIDE 18: Installation & Deployment (Light Theme)
    # ==========================================
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18, c_light_bg)
    add_slide_header(slide18, "Installation & Deployment Protocols", "Streamlined Setup Options for Production and Testing")
    
    # Left Column: Option 1 - Docker Compose
    draw_card(slide18, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide18.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Option 1: Docker Compose Setup", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf, "Recommended for quick starting multi-container stacks:", font_size=12, space_after=8)
    
    # Command blocks representation
    draw_card(slide18, Inches(1.0), Inches(3.1), Inches(5.2), Inches(1.3), bg_color=c_dark_bg, border_color=c_dark_bg)
    tx_cmd = slide18.shapes.add_textbox(Inches(1.1), Inches(3.2), Inches(5.0), Inches(1.1))
    tf_cmd = tx_cmd.text_frame
    add_p(tf_cmd, "$ docker-compose up -d", font_size=11, font_color=c_white)
    add_p(tf_cmd, "$ docker exec -it ollama-service ollama pull llama3", font_size=11, font_color=c_white)
    add_p(tf_cmd, "Access: http://localhost:8501", font_size=11, font_color=c_accent_teal, bold=True)
    
    add_p(tf, "Compose Configuration advantages:", font_size=12, bold=True, space_after=4)
    add_p(tf, "• Port Binding: Deploys Ollama on port 11434 and Streamlit on port 8501.", font_size=10, font_color=c_muted_text, space_after=4)
    add_p(tf, "• Socket Mounting: Mounts docker.sock directly with read-only (:ro) permissions for safety.", font_size=10, font_color=c_muted_text)

    # Right Column: Option 2 - Manual Setup
    draw_card(slide18, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_man = slide18.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_man = tx_man.text_frame
    tf_man.word_wrap = True
    add_p(tf_man, "Option 2: Standalone Shell Setup", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf_man, "Recommended for manual local development execution:", font_size=12, space_after=8)
    
    # Command blocks representation
    draw_card(slide18, Inches(7.1), Inches(3.1), Inches(5.2), Inches(1.3), bg_color=c_dark_bg, border_color=c_dark_bg)
    tx_cmd2 = slide18.shapes.add_textbox(Inches(7.2), Inches(3.2), Inches(5.0), Inches(1.1))
    tf_cmd2 = tx_cmd2.text_frame
    add_p(tf_cmd2, "$ pip install -r requirements.txt", font_size=11, font_color=c_white)
    add_p(tf_cmd2, "$ ollama run llama3", font_size=11, font_color=c_white)
    add_p(tf_cmd2, "$ streamlit run app.py", font_size=11, font_color=c_white)
    
    add_p(tf_man, "Environmental Configurations:", font_size=12, bold=True, space_after=4)
    add_p(tf_man, "• OLLAMA_HOST: Configures remote Ollama endpoints (defaulting to http://localhost:11434).", font_size=10, font_color=c_muted_text, space_after=4)
    add_p(tf_man, "• DOCKER_HOST: Maps custom daemon endpoints when running Docker remotely.", font_size=10, font_color=c_muted_text)

    # ==========================================
    # SLIDE 19: Troubleshooting & QA checklists (Light Theme)
    # ==========================================
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19, c_light_bg)
    add_slide_header(slide19, "Troubleshooting & QA Checklists", "Resolving environment conflicts and running testing scripts")
    
    # Left Column: Troubleshooting
    draw_card(slide19, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.6))
    tx = slide19.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tx.text_frame
    tf.word_wrap = True
    add_p(tf, "Handling Common Failures", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    
    add_p(tf, "1. Docker Connection Socket Failures", font_size=12, bold=True, space_after=2)
    add_p(tf, "• Reason: Current user lacks permissions to bind to docker.sock.\n• Fix: Linux user execution command: sudo usermod -aG docker $USER.", font_size=10, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "2. Ollama Connection Timeout Failures", font_size=12, bold=True, space_after=2)
    add_p(tf, "• Reason: Ollama backend service is shut down or model is missing.\n• Fix: Confirm service run via: curl http://localhost:11434/api/tags. Pull model: ollama pull llama3.", font_size=10, font_color=c_muted_text, space_after=8)
    
    add_p(tf, "3. UI Port Conflicted Failures", font_size=12, bold=True, space_after=2)
    add_p(tf, "• Reason: Streamlit port 8501 is locked by a zombie task.\n• Fix: Kill using: lsof -ti:8501 | xargs kill -9. Or: streamlit run app.py --server.port 8502.", font_size=10, font_color=c_muted_text)

    # Right Column: QA scripts
    draw_card(slide19, Inches(6.9), Inches(2.1), Inches(5.6), Inches(4.6))
    tx_qa = slide19.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_qa = tx_qa.text_frame
    tf_qa.word_wrap = True
    add_p(tf_qa, "Automated Verification Testing", font_size=16, font_color=c_accent_blue, bold=True, space_after=12)
    add_p(tf_qa, "The package includes test_demo.py to automate core functional checks in staging:", font_size=11, space_after=8)
    
    add_p(tf_qa, "• Mock Container Lifecycle creation", font_size=12, bold=True, space_after=2)
    add_p(tf_qa, "Auto-deploys 3 testing containers: nginx-test, redis-test, mysql-test.", font_size=10, font_color=c_muted_text, space_after=6)
    
    add_p(tf_qa, "• Natural Language Prompt simulations", font_size=12, bold=True, space_after=2)
    add_p(tf_qa, "Simulates standard actions: restarting, stopping, checking resource limits.", font_size=10, font_color=c_muted_text, space_after=6)
    
    add_p(tf_qa, "• Test Execution commands", font_size=12, bold=True, space_after=2)
    add_p(tf_qa, "Run verification test scripts: python test_demo.py.", font_size=10, font_color=c_muted_text)

    # ==========================================
    # SLIDE 20: Future Scope & Conclusion (Dark Theme)
    # ==========================================
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20, c_dark_bg)
    
    accent_bar = slide20.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_accent_teal
    accent_bar.line.color.rgb = c_accent_teal
    
    # Title
    title_box = slide20.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(1.2))
    tf20 = title_box.text_frame
    tf20.word_wrap = True
    p20 = tf20.paragraphs[0]
    p20.text = "Roadmap & Conclusion"
    p20.font.name = f_sans
    p20.font.size = Pt(36)
    p20.font.bold = True
    p20.font.color.rgb = c_white
    
    # Text block with lists of future roadmap
    tx_road = slide20.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(11.0), Inches(4.0))
    tf_road = tx_road.text_frame
    tf_road.word_wrap = True
    add_p(tf_road, "Future Engineering Scope:", font_size=16, font_color=c_accent_teal, bold=True, space_after=10)
    add_p(tf_road, "• Remote Multihost Controls: Extend dashboard connection configuration to handle cluster nodes.", font_size=12, font_color=c_white, space_after=6)
    add_p(tf_road, "• Docker Compose management: Add agent rules to parse, configure, and boot docker-compose files.", font_size=12, font_color=c_white, space_after=6)
    add_p(tf_road, "• Live Telemetry Sockets: Replace pandas grid intervals with asynchronous WebSockets streams.", font_size=12, font_color=c_white, space_after=6)
    add_p(tf_road, "• Granular Access controls: Implement role-based credentials (RBAC) separating viewers from executioners.", font_size=12, font_color=c_white, space_after=20)
    
    add_p(tf_road, "Key Takeaways:", font_size=16, font_color=c_accent_teal, bold=True, space_after=10)
    add_p(tf_road, "✔ Bridges the command line skills gap safely via controlled, schema-locked ReAct loops.", font_size=12, font_color=c_white, space_after=6)
    add_p(tf_road, "✔ Enforces zero external API exposure by executing translation models locally on premise.", font_size=12, font_color=c_white, space_after=6)
    add_p(tf_road, "✔ Full telemetry visual tracking via unified historical database records and interactive charts.", font_size=12, font_color=c_white, space_after=12)

    # Save presentation
    prs.save("Docker_NL_Dashboard_Presentation.pptx")
    prs.save("Docker_NL_Dashboard_Presentation_v3.pptx")
    print("Presentation compiled successfully with exactly 20 slides!")

if __name__ == "__main__":
    create_presentation()
